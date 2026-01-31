import os
import pandas as pd
from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader, CSVLoader
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import OllamaEmbeddings
from langchain_community.vectorstores import Milvus
from langchain_core.documents import Document
from app.core.config import config

def ingest_documents():
    data_dir = "data"
    if not os.path.exists(data_dir):
        os.makedirs(data_dir)
        
    all_documents = []
    
    # scan directory
    for filename in os.listdir(data_dir):
        file_path = os.path.join(data_dir, filename)
        
        # Process PDF
        if filename.endswith(".pdf"):
            print(f"Loading PDF: {filename}")
            try:
                loader = PyPDFLoader(file_path)
                docs = loader.load()
                # Initial split
                text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
                splits = text_splitter.split_documents(docs)
                for split in splits:
                    split.metadata["source"] = filename
                    split.metadata["type"] = "pdf"
                all_documents.extend(splits)
            except Exception as e:
                print(f"Error loading PDF {filename}: {e}")

        # Process Excel
        elif filename.endswith(".xlsx") or filename.endswith(".xls"):
            print(f"Loading Excel: {filename}")
            try:
                df = pd.read_excel(file_path)
                # Convert rows to documents
                for index, row in df.iterrows():
                    # Create a text representation of the row
                    row_text = ", ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                    doc = Document(
                        page_content=row_text,
                        metadata={
                            "source": filename,
                            "row": index,
                            "type": "excel"
                        }
                    )
                    all_documents.append(doc)
            except Exception as e:
                print(f"Error loading Excel {filename}: {e}")

        # Process Text Files
        elif filename.endswith(".txt"):
            print(f"Loading Text File: {filename}")
            try:
                loader = TextLoader(file_path, encoding="utf-8")
                docs = loader.load()
                # Split text
                text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
                splits = text_splitter.split_documents(docs)
                for split in splits:
                    split.metadata["source"] = filename
                    split.metadata["type"] = "text"
                all_documents.extend(splits)
            except Exception as e:
                print(f"Error loading Text File {filename}: {e}")

        # Process CSV
        elif filename.endswith(".csv"):
            print(f"Loading CSV: {filename}")
            try:
                loader = CSVLoader(file_path=file_path)
                docs = loader.load()
                # CSV loader creates one doc per row usually, so we might just extend
                for doc in docs:
                    doc.metadata["source"] = filename
                    doc.metadata["type"] = "csv"
                all_documents.extend(docs)
            except Exception as e:
                print(f"Error loading CSV {filename}: {e}")

        # Process Word (DOCX)
        elif filename.endswith(".docx") or filename.endswith(".doc"):
            print(f"Loading Word Doc: {filename}")
            try:
                loader = Docx2txtLoader(file_path)
                docs = loader.load()
                text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
                splits = text_splitter.split_documents(docs)
                for split in splits:
                    split.metadata["source"] = filename
                    split.metadata["type"] = "docx"
                all_documents.extend(splits)
            except Exception as e:
                print(f"Error loading DOCX {filename}: {e}")

        # Process PowerPoint (PPTX)
        elif filename.endswith(".pptx") or filename.endswith(".ppt"):
            print(f"Loading PowerPoint: {filename}")
            try:
                prs = Presentation(file_path)
                text_content = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content += shape.text + "\n"
                
                doc = Document(page_content=text_content, metadata={"source": filename, "type": "pptx"})
                
                text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
                splits = text_splitter.split_documents([doc])
                all_documents.extend(splits)
            except Exception as e:
                print(f"Error loading PPTX {filename}: {e}")

    if not all_documents:
        print("No documents found to ingest!")
        return

    print(f"Total documents to ingest: {len(all_documents)}")
    
    print("Initializing Embeddings (nomic-embed-text)...")
    embeddings = OllamaEmbeddings(
        model=config.EMBEDDING_MODEL,
        base_url=config.OLLAMA_BASE_URL
    )
    
    # Sanitize metadata for Milvus (Auto-schema prefers consistent types)
    for doc in all_documents:
        new_metadata = {}
        for k, v in doc.metadata.items():
            if isinstance(v, (str, int, float, bool)):
                new_metadata[k] = str(v)  # Convert everything to string for safety
        doc.metadata = new_metadata

    print(f"Indexing to Milvus collection '{config.COLLECTION_NAME}'...")
    try:
        Milvus.from_documents(
            all_documents,
            embeddings,
            collection_name=config.COLLECTION_NAME,
            connection_args={"host": config.MILVUS_HOST, "port": config.MILVUS_PORT},
            drop_old=True # Reset collection for fresh start
        )
        print("Indexing to Milvus Complete!")
        
        # Save chunks for BM25 (Hybrid Search)
        import pickle
        with open(os.path.join(data_dir, "chunks.pkl"), "wb") as f:
            pickle.dump(all_documents, f)
        print("Saved chunks for Hybrid Search.")
        
    except Exception as e:
        print(f"Failed to ingest to Milvus: {e}")

if __name__ == "__main__":
    ingest_documents()
